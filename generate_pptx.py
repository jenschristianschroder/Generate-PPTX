"""
Generate a PowerPoint from Microsoft Dataverse data using a fixed template.

Dependencies:
    pip install msal requests python-pptx

Configuration:
    Set the following environment variables:
        DATAVERSE_CLIENT_ID   - Azure AD application (client) ID
        DATAVERSE_CLIENT_SECRET - Azure AD application client secret
        DATAVERSE_TENANT_ID   - Azure AD tenant ID
        DATAVERSE_URL         - Dataverse endpoint, e.g. https://<org>.api.crm.dynamics.com/api/data/v9.1/

Usage:
    python generate_ppt_dataverse.py
"""
import os
import sys
import msal
import requests
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from dotenv import load_dotenv
from datetime import datetime

# Load environment variables from a .env file
load_dotenv()

CLIENT_ID = os.getenv('DATAVERSE_CLIENT_ID')
CLIENT_SECRET = os.getenv('DATAVERSE_CLIENT_SECRET')
TENANT_ID = os.getenv('DATAVERSE_TENANT_ID')
DATAVERSE_URL = os.getenv('DATAVERSE_URL')
DATAVERSE_API_URL = os.getenv('DATAVERSE_API_URL')

DATAVERSE_ENTITY = os.getenv('DATAVERSE_ENTITY')
DATAVERSE_ENTITY_COLUMNS = os.getenv('DATAVERSE_ENTITY_COLUMNS')
DATAVERSE_ENTITY_FILTER_COLUMN = os.getenv('DATAVERSE_ENTITY_FILTER_COLUMN')

PPTX_TEMPLATE = os.getenv('PPTX_TEMPLATE', 'template.pptx')
PPTX_TABLE_STYLE_ID = os.getenv('PPTX_TABLE_STYLE_ID')

base_path = os.getcwd()
OUTPUT_PATH = os.path.join(base_path, "output")

if not all([CLIENT_ID, CLIENT_SECRET, TENANT_ID, DATAVERSE_API_URL, DATAVERSE_URL]):
    print("Please set DATAVERSE_CLIENT_ID, DATAVERSE_CLIENT_SECRET, DATAVERSE_TENANT_ID, DATAVERSE_API_URL and DATAVERSE_URL.")
    sys.exit(1)

SCOPE = [f"{DATAVERSE_URL.strip('/')}/.default"]


def get_access_token():
    """Retrieve an OAuth2 access token for Dataverse."""
    authority = f"https://login.microsoftonline.com/{TENANT_ID}"
    app = msal.ConfidentialClientApplication(
        CLIENT_ID,
        authority=authority,
        client_credential=CLIENT_SECRET
    )
    result = app.acquire_token_for_client(scopes=SCOPE)
    if 'access_token' in result:
        return result['access_token']
    else:
        raise Exception(f"Token acquisition failed: {result.get('error_description')}")


def fetch_data(entity: str, token: str, select=None, filter_expr=None):
    """Retrieve records from a Dataverse entity."""
    headers = {
        'Authorization': f'Bearer {token}',
        'Accept': 'application/json'
    }
    params = {}
    if select:
        params['$select'] = ','.join(select)
    if filter_expr:
        params['$filter'] = filter_expr

    url = f"{DATAVERSE_API_URL.rstrip('/')}/{entity}"
    resp = requests.get(url, headers=headers, params=params)
    resp.raise_for_status()
    return resp.json().get('value', [])


def iter_cells(table):
    """Helper function to iterate over all cells in a table."""
    for row in table.rows:
        for cell in row.cells:
            yield cell


def set_table_font_size(table, font_size):
    """Set the font size for all cells in a table."""
    for cell in iter_cells(table):
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(font_size)


def create_table(slide, rows, cols, left, top, width, height, font_size, font_bold):
    """Create a new table and apply formatting."""
    new_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    new_table = new_shape.table
    tbl =  new_shape._element.graphic.graphicData.tbl
    tbl[0][-1].text = PPTX_TABLE_STYLE_ID

    # Apply font size and bold formatting to all cells
    for cell in iter_cells(new_table):
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = font_size
                run.font.bold = font_bold

    return new_table


def process_table_placeholder(table, content, slide):
    """Process a table placeholder and replace it with a new table."""
    placeholder_name = table.cell(0, 0).text[8:-2].strip()  # Extract placeholder name
    #print(f"Found table placeholder: {placeholder_name}")

    # Get the array from the JSON content
    value = content.get(placeholder_name, [])
    if not isinstance(value, list) or not value:
        print(f"No valid array found for '{placeholder_name}'")
        table.cell(0, 0).text = "n/a"
        set_table_font_size(table, 11)
        return

    # Determine headers and table dimensions
    if isinstance(value[0], dict):
        headers = list(value[0].keys())
        rows, cols = len(value) + 1, len(headers)  # +1 for the header row
    else:
        print(f"Invalid array format for '{placeholder_name}'")
        return

    # Extract formatting and position of the placeholder table
    left = table._graphic_frame.left
    top = table._graphic_frame.top
    width = table._graphic_frame.width
    height = table._graphic_frame.height
    font_size = table.cell(0, 0).text_frame.paragraphs[0].font.size
    font_bold = table.cell(0, 0).text_frame.paragraphs[0].font.bold


    # Remove the placeholder table
    sp = table._graphic_frame._element
    sp.getparent().remove(sp)

    # Create a new table
    new_table = create_table(slide, rows, cols, left, top, width, height, font_size, font_bold)

    # Insert headers
    for col_idx, header in enumerate(headers):
        cell = new_table.cell(0, col_idx)
        cell.text = header

    # Insert data rows
    for row_idx, item in enumerate(value, start=1):
        for col_idx, header in enumerate(headers):
            cell = new_table.cell(row_idx, col_idx)
            cell.text = str(item.get(header, ""))

    # Set font size for all cells
    set_table_font_size(new_table, 11)


def process_text_placeholders(slide, content):
    """
    Update text placeholders in a slide.
    Replace {{placeholders}} in text elements with corresponding values from the JSON content.
    If a placeholder is not found in the JSON document, replace it with "n/a".
    Args:
        slide: The slide object to process.
        content: The JSON content containing placeholder values.
    """
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                # Concatenate all runs into a single string
                fulltext = ''.join(run.text for run in paragraph.runs)

                # Replace placeholders in the concatenated text
                for placeholder_name, value in content.items():
                    placeholder_tag = f"{{{{{placeholder_name}}}}}"  # Format as {{placeholder}}
                    if placeholder_tag in fulltext:
                        #print(f"Replacing placeholder: {placeholder_name}")
                        fulltext = fulltext.replace(placeholder_tag, str(value))

                # Handle placeholders not found in the JSON document
                while "{{" in fulltext and "}}" in fulltext:
                    start_idx = fulltext.find("{{")
                    end_idx = fulltext.find("}}", start_idx) + 2
                    placeholder_tag = fulltext[start_idx:end_idx]
                    #print(f"Replacing missing placeholder: {placeholder_tag} with 'n/a'")
                    fulltext = fulltext.replace(placeholder_tag, "n/a")

                # Clear all runs and reassign the updated text
                for run in paragraph.runs:
                    run.text = ''  # Clear existing text
                if paragraph.runs:
                    paragraph.runs[0].text = fulltext  # Assign updated text to the first run


def generate_ppt(jobid: str, records: list, template_path='template.pptx', output_filename='output.pptx'):
    """Populate a PowerPoint template with data records."""
    prs = Presentation(template_path)
    jobdate = datetime.now().strftime('%Y-%m-%d %H:%M:%S')

    for idx, record in enumerate(records):
        content = {}
        if 'jeschro_content' in record:
            try:
                content = json.loads(record['jeschro_content'])
            except json.JSONDecodeError as e:
                print(f"Error decoding JSON for record {idx}: {e}")
                continue

        content['jobid'] = jobid
        content['jobdate'] = jobdate

        for slide in prs.slides:
            # Update text placeholders
            process_text_placeholders(slide, content)

            # Process table placeholders
            for shape in slide.shapes:
                if shape.has_table:
                    process_table_placeholder(shape.table, content, slide)

    print(OUTPUT_PATH)    
    if not os.path.exists(OUTPUT_PATH):
        os.makedirs(OUTPUT_PATH)

    prs.save(os.path.join(OUTPUT_PATH, output_filename))
    print(f"Generated presentation: {os.path.join(OUTPUT_PATH, output_filename)}")


if __name__ == '__main__':
    token = get_access_token()

    #jobid = "d13b4413-f120-f011-9989-7c1e5283aeb9"
    jobid = "10960a97-4621-f011-8c4d-7c1e5283aeb9"
   
    data = fetch_data(
        entity=DATAVERSE_ENTITY,
        token=token,
        select=[DATAVERSE_ENTITY_COLUMNS],
        filter_expr="".join([DATAVERSE_ENTITY_FILTER_COLUMN, f" eq '{jobid}'"])
    )

    if not data:
        print("No records retrieved.")
        sys.exit(0)

    generate_ppt(
        jobid=jobid,
        records=data,
        template_path=PPTX_TEMPLATE,
        output_filename='dataverse_report.pptx'
    )
